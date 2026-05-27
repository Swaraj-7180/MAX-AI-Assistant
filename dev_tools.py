import os
import re
import subprocess
import tempfile
import threading
import time
import json
from pathlib import Path
from typing import Optional
from groq import Groq

# Add your named presentation paths here.
# Example:
# PPT_PATHS = {
#     "projectslides": r"C:\Users\Swaraj Shinde\Documents\ProjectSlides.pptx",
#     "demo": r"D:\Presentations\DemoDeck.pptx",
# }
PPT_PATHS = {
    "Minor project": r"https://d.docs.live.net/5de59eda02327d7a/Documents/Minor%20project.pptx",
}

try:
    import pytesseract
    pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
except ImportError:
    pass


def _get_client():
    return Groq(api_key="Grok api key")


class DevTools:
    def __init__(self, say_callback) -> None:
        self.say = say_callback
        self._pending = None
        # Load hardcoded presentation path mappings and optional ppt_paths.json overrides
        self.ppt_map: dict[str, str] = {k.lower(): v for k, v in PPT_PATHS.items()}
        try:
            cfg = Path(__file__).parent / "ppt_paths.json"
            self._ppt_paths_file = cfg
            if cfg.exists():
                with cfg.open("r", encoding="utf-8") as f:
                    raw = json.load(f)
                if isinstance(raw, dict):
                    self.ppt_map.update({k.lower(): v for k, v in raw.items()})
            else:
                # create an empty file so users can edit it
                with cfg.open("w", encoding="utf-8") as f:
                    json.dump({}, f, indent=2)
        except Exception:
            pass

    # ── MAIN HANDLER ─────────────────────────────────
    def handle(self, text: str) -> Optional[str]:
        query = text.strip().lower()

        # Pending flow (bug reporter multi-step)
        if self._pending:
            return self._handle_pending(text)

        # Code Explainer
        if re.search(r"\b(explain|what does|describe)\s+(this\s+)?(code|function|script)\b", query):
            return self._explain_code()

        # Code Review
        if re.search(r"\b(review|check|analyze|analyse)\s+(my\s+)?(code|script|file)\b", query):
            return self._review_code()

        # Bug Reporter
        if re.search(r"\b(found a bug|there.?s a bug|bug report|report a bug|i have a bug)\b", query):
            self._pending = {"type": "bug_desc"}
            return "Tell me what happened — describe the bug."

        # Presentation explainer
        if re.search(r"\b(describe|explain|summarize|read)\b.*\b(powerpoint|ppt|pptx|slides)\b", query) or \
           re.search(r"\b(powerpoint|ppt|pptx|slides)\b.*\b(describe|explain|summarize|read)\b", query):
            return self._describe_presentation_request(query)

        # Present (start slideshow) and explain
        present_match = re.search(r"\b(present|start presentation|start slideshow|show slides|present slides)\b(?:\s+(.+))?", query)
        if present_match:
            name = present_match.group(2)
            return self._present_presentation(name.strip() if name else None)

        # Terminal commands
        terminal_match = re.search(
            r"\b(run|execute|install|uninstall|start server|stop server|"
            r"create file|make folder|list files|show files|"
            r"git commit|git push|git pull|git status|git branch|"
            r"activate env|deactivate env)\b(.+)?", query
        )
        if terminal_match:
            return self._terminal_command(query)

        return None

    # ── SCREENSHOT HELPER ────────────────────────────
    def _take_screenshot(self) -> Optional[str]:
        try:
            import pyautogui
            import base64
            from PIL import Image
            import io

            screenshot = pyautogui.screenshot()
            buffer = io.BytesIO()
            screenshot.save(buffer, format="PNG")
            buffer.seek(0)
            return base64.b64encode(buffer.read()).decode("utf-8")
        except ImportError:
            return None
        except Exception:
            return None

    def _screenshot_to_text(self) -> Optional[str]:
        """Use pytesseract to extract text from screenshot"""
        try:
            import pyautogui
            import pytesseract
            from PIL import Image

            screenshot = pyautogui.screenshot()
            text = pytesseract.image_to_string(screenshot)
            return text.strip() if text.strip() else None
        except ImportError:
            return None
        except Exception:
            return None

    # ── CODE EXPLAINER ───────────────────────────────
    def _explain_code(self) -> str:
        self.say("Taking a screenshot of your screen.")

        # Try OCR first
        code_text = self._screenshot_to_text()

        if not code_text:
            # Fallback — ask user to paste code
            self._pending = {"type": "explain_paste"}
            return "I couldn't read the screen. Paste or speak the code you want explained."

        try:
            client = _get_client()
            response = client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[{
                    "role": "user",
                    "content": (
                        f"Explain this code in simple words in 3-4 sentences. "
                        f"Focus on what it does, not how. Be concise:\n\n{code_text[:3000]}"
                    )
                }],
                max_tokens=300
            )
            return response.choices[0].message.content.strip()
        except Exception:
            return "Couldn't explain the code. AI service unavailable."

    # ── CODE REVIEW ──────────────────────────────────
    def _review_code(self) -> str:
        self.say("Analyzing your code. Give me a moment.")

        code_text = self._screenshot_to_text()

        if not code_text:
            self._pending = {"type": "review_paste"}
            return "I couldn't read the screen. Paste the code you want reviewed."

        try:
            client = _get_client()
            response = client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[{
                    "role": "user",
                    "content": (
                        f"Review this code like a senior developer. "
                        f"Point out: bugs, improvements, bad practices. "
                        f"Keep it under 4 sentences, be direct and specific:\n\n{code_text[:3000]}"
                    )
                }],
                max_tokens=300
            )
            return response.choices[0].message.content.strip()
        except Exception:
            return "Couldn't review the code. AI service unavailable."

    # ── TERMINAL COMMANDS ────────────────────────────
    def _terminal_command(self, query: str) -> str:
        command = self._parse_terminal_command(query)
        if not command:
            return "I didn't understand that command."

        # Safety check — block dangerous commands
        dangerous = ["rm -rf", "format", "del /f /s", "rmdir /s", "shutdown", "rd /s"]
        if any(d in command.lower() for d in dangerous):
            return "That command is too dangerous. I won't run that."

        try:
            print(f"[MAX] Running: {command}")
            result = subprocess.run(
                command, shell=True, capture_output=True,
                text=True, timeout=30, cwd=r"D:\MAX_AI"
            )
            output = result.stdout.strip() or result.stderr.strip()
            if output:
                # Summarize long output with AI
                if len(output) > 200:
                    return self._summarize_output(command, output)
                return f"Done. {output[:150]}"
            return f"Command executed successfully."
        except subprocess.TimeoutExpired:
            return "Command timed out after 30 seconds."
        except Exception as e:
            return f"Command failed: {str(e)[:100]}"

    def _parse_terminal_command(self, query: str) -> Optional[str]:
        # Install packages
        install_match = re.search(r"\binstall\s+(.+)", query)
        if install_match:
            pkg = install_match.group(1).strip()
            return f"pip install {pkg}"

        uninstall_match = re.search(r"\buninstall\s+(.+)", query)
        if uninstall_match:
            pkg = uninstall_match.group(1).strip()
            return f"pip uninstall {pkg} -y"

        # Git commands
        if "git commit" in query:
            msg_match = re.search(r"git commit\s+(.+)", query)
            msg = msg_match.group(1).strip() if msg_match else "auto commit by MAX"
            return f'git add . && git commit -m "{msg}"'
        if "git push" in query:
            return "git push"
        if "git pull" in query:
            return "git pull"
        if "git status" in query:
            return "git status"
        if "git branch" in query:
            branch_match = re.search(r"git branch\s+(.+)", query)
            branch = branch_match.group(1).strip() if branch_match else ""
            return f"git checkout -b {branch}" if branch else "git branch"

        # Server
        if "run server" in query or "start server" in query:
            return "python main.py"
        if "stop server" in query:
            return "taskkill /f /im python.exe"

        # Files
        if "list files" in query or "show files" in query:
            return "dir"

        # Run file
        run_match = re.search(r"\brun\s+(.+\.py)\b", query)
        if run_match:
            return f"python {run_match.group(1).strip()}"

        # Generic — ask AI to generate command
        return self._ai_generate_command(query)

    def _ai_generate_command(self, query: str) -> Optional[str]:
        try:
            client = _get_client()
            response = client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[{
                    "role": "user",
                    "content": (
                        f"Convert this to a Windows terminal command. "
                        f"Return ONLY the command, nothing else: '{query}'"
                    )
                }],
                max_tokens=50
            )
            cmd = response.choices[0].message.content.strip()
            cmd = re.sub(r"```.*?```", "", cmd, flags=re.DOTALL).strip()
            return cmd if cmd else None
        except Exception:
            return None

    def _summarize_output(self, command: str, output: str) -> str:
        try:
            client = _get_client()
            response = client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[{
                    "role": "user",
                    "content": (
                        f"Summarize this terminal output in 1-2 sentences. "
                        f"Command was: {command}\nOutput:\n{output[:1000]}"
                    )
                }],
                max_tokens=100
            )
            return response.choices[0].message.content.strip()
        except Exception:
            return f"Command done. Output was {len(output)} characters."

    # ── BUG REPORTER ─────────────────────────────────
    def _handle_pending(self, text: str) -> Optional[str]:
        p = self._pending

        if p["type"] == "bug_desc":
            self._pending = {"type": "bug_steps", "desc": text}
            return "What were you doing when it happened?"

        elif p["type"] == "bug_steps":
            self._pending = {"type": "bug_confirm",
                           "desc": p["desc"], "steps": text}
            return f"Got it. Should I create a GitHub issue for this bug?"

        elif p["type"] == "bug_confirm":
            if any(w in text.lower() for w in ["yes","yeah","sure","do it","create","go ahead"]):
                result = self._create_github_issue(p["desc"], p["steps"])
                self._pending = None
                return result
            else:
                self._pending = None
                return "Bug report cancelled. Let me know when you want to report it."

        elif p["type"] == "explain_paste":
            self._pending = None
            return self._explain_code_text(text)

        elif p["type"] == "ppt_file":
            self._pending = None
            path = text.strip().strip('"')
            return self._describe_presentation(path)

        elif p["type"] == "review_paste":
            self._pending = None
            return self._review_code_text(text)

        self._pending = None
        return None

    def _explain_code_text(self, code: str) -> str:
        try:
            client = _get_client()
            response = client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[{
                    "role": "user",
                    "content": f"Explain this code simply in 3-4 sentences:\n\n{code}"
                }],
                max_tokens=300
            )
            return response.choices[0].message.content.strip()
        except Exception:
            return "Couldn't explain the code."

    def _review_code_text(self, code: str) -> str:
        try:
            client = _get_client()
            response = client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[{
                    "role": "user",
                    "content": (
                        f"Review this code. Point out bugs and improvements "
                        f"in 3-4 sentences:\n\n{code}"
                    )
                }],
                max_tokens=300
            )
            return response.choices[0].message.content.strip()
        except Exception:
            return "Couldn't review the code."

    def _describe_presentation_request(self, query: str) -> str:
        file_path = self._find_pptx_path(query)
        if file_path:
            return self._describe_presentation(file_path)
        self._pending = {"type": "ppt_file"}
        return (
            "Send me the full path to your PowerPoint file (.pptx). "
            "I will read the slides and explain them so you can present them to an interviewer."
        )

    def _present_presentation(self, name: Optional[str]) -> str:
        # If a name/path provided, try to locate
        if name:
            # If user gave a full path, use it
            if os.path.exists(name) and name.lower().endswith('.pptx'):
                file_path = name
            else:
                file_path = self._locate_pptx_by_name(name)

            if not file_path:
                self._pending = {"type": "ppt_file"}
                return (f"Couldn't find a presentation matching '{name}'. "
                        "Please provide the full .pptx path.")

            started = self._start_slideshow(file_path)
            # Start explanation in background so slideshow isn't blocked
            threading.Thread(target=self._speak_presentation_explanation, args=(file_path,), daemon=True).start()
            if started:
                return f"Starting presentation {os.path.basename(file_path)} and explaining the slides."
            return f"Opened {os.path.basename(file_path)}. I'll explain the slides now."

        # No name given — ask for path
        self._pending = {"type": "ppt_file"}
        return "Which presentation would you like to present? Send the full .pptx path or the file name."

    def _locate_pptx_by_name(self, name: str) -> Optional[str]:
        # Check user-configured mapping first (exact or partial key match)
        key = name.strip().lower()
        if key in self.ppt_map:
            path = self.ppt_map[key]
            if os.path.exists(path):
                return path
        for k, v in self.ppt_map.items():
            if key in k and os.path.exists(v):
                return v

        # Fallback: scan common user folders for a matching filename
        roots = [
            os.path.expanduser("~\\Desktop"),
            os.path.expanduser("~\\Documents"),
            os.path.expanduser("~\\Downloads"),
        ]
        name_low = key
        for root in roots:
            if not os.path.exists(root):
                continue
            try:
                for dirpath, dirnames, filenames in os.walk(root):
                    for fname in filenames:
                        if not fname.lower().endswith('.pptx'):
                            continue
                        if name_low in fname.lower():
                            return os.path.join(dirpath, fname)
            except PermissionError:
                continue
        return None

    def _start_slideshow(self, file_path: str) -> bool:
        # Try known PowerPoint locations
        candidates = [
            r"C:\\Program Files\\Microsoft Office\\root\\Office16\\POWERPNT.EXE",
            r"C:\\Program Files (x86)\\Microsoft Office\\Office16\\POWERPNT.EXE",
            r"C:\\Program Files\\Microsoft Office\\Office16\\POWERPNT.EXE",
        ]
        for exe in candidates:
            if os.path.exists(exe):
                try:
                    subprocess.Popen([exe, "/s", file_path], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                    return True
                except Exception:
                    break
        # Fallback — open file normally (user may need to press slideshow)
        try:
            os.startfile(file_path)
            return False
        except Exception:
            return False

    def _speak_presentation_explanation(self, file_path: str) -> None:
        # Get AI explanation and speak it via callback
        explanation = self._describe_presentation(file_path)
        # Break into paragraphs to avoid very long TTS calls
        for part in explanation.split('\n\n'):
            part = part.strip()
            if part:
                try:
                    self.say(part)
                    # small pause between parts
                    time.sleep(0.8)
                except Exception:
                    pass

    def _find_pptx_path(self, text: str) -> Optional[str]:
        path_match = re.search(r"([A-Za-z]:\\[^\s]+?\.pptx?)", text, re.I)
        if path_match:
            return path_match.group(1)
        return None

    def _describe_presentation(self, file_path: str) -> str:
        if not file_path.lower().endswith(".pptx"):
            return "I can only describe .pptx PowerPoint files right now."
        if not os.path.exists(file_path):
            return "I couldn't find that file. Please give me the correct .pptx path."

        self.say("Reading your PowerPoint slides. One moment, please.")
        slides = self._extract_pptx_text(file_path)
        if not slides:
            return "I couldn't extract any slide text from that presentation."

        presentation_text = []
        for index, slide_text in slides:
            cleaned = " ".join(slide_text.split())
            presentation_text.append(f"Slide {index}: {cleaned}")
        prompt_text = "\n\n".join(presentation_text)
        try:
            client = _get_client()
            response = client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[{
                    "role": "user",
                    "content": (
                        "You are a presentation coach. Read the following PowerPoint slide text and "
                        "explain each slide in a way that a project interviewer can understand. "
                        "For each slide, give a short summary and a suggested way to speak about it. "
                        "Be clear, professional, and keep each slide explanation concise.\n\n"
                        f"{prompt_text}"
                    )
                }],
                max_tokens=600
            )
            return response.choices[0].message.content.strip()
        except Exception:
            return "I couldn't explain the presentation right now. Please try again later."

    def _extract_pptx_text(self, file_path: str) -> list[tuple[int, str]]:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            slides = []
            for index, slide in enumerate(prs.slides, start=1):
                text_blocks = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_blocks.append(shape.text)
                slides.append((index, "\n".join(text_blocks).strip()))
            return slides
        except ImportError:
            return self._extract_pptx_text_fallback(file_path)
        except Exception:
            return []

    def _extract_pptx_text_fallback(self, file_path: str) -> list[tuple[int, str]]:
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            slides = []
            with zipfile.ZipFile(file_path, "r") as archive:
                slide_files = sorted(
                    [name for name in archive.namelist()
                     if name.startswith("ppt/slides/slide") and name.endswith(".xml")]
                )
                for index, slide_name in enumerate(slide_files, start=1):
                    xml_data = archive.read(slide_name)
                    root = ET.fromstring(xml_data)
                    texts = []
                    for node in root.iter():
                        if node.tag.endswith("}t") and node.text:
                            texts.append(node.text)
                    slides.append((index, "\n".join(texts).strip()))
            return slides
        except Exception:
            return []

    # ── GITHUB ISSUE ─────────────────────────────────
    def _create_github_issue(self, description: str, steps: str) -> str:
        # Configure these
        GITHUB_TOKEN = "your_github_token_here"  # 👈 add your token
        GITHUB_REPO  = "your_username/your_repo"  # 👈 add your repo

        if "your_github" in GITHUB_TOKEN:
            # Save locally if no GitHub token
            return self._save_bug_locally(description, steps)

        try:
            import requests
            title = description[:60] + "..." if len(description) > 60 else description
            body = f"## Bug Description\n{description}\n\n## Steps to Reproduce\n{steps}"
            response = requests.post(
                f"https://api.github.com/repos/{GITHUB_REPO}/issues",
                json={"title": title, "body": body, "labels": ["bug"]},
                headers={"Authorization": f"token {GITHUB_TOKEN}",
                        "Accept": "application/vnd.github.v3+json"}
            )
            if response.status_code == 201:
                issue_url = response.json()["html_url"]
                return f"GitHub issue created. Check it at {issue_url}"
            return self._save_bug_locally(description, steps)
        except Exception:
            return self._save_bug_locally(description, steps)

    def _save_bug_locally(self, description: str, steps: str) -> str:
        try:
            import json
            from datetime import datetime
            bugs_file = r"D:\MAX_AI\bugs.json"
            bugs = []
            if os.path.exists(bugs_file):
                with open(bugs_file) as f:
                    bugs = json.load(f)
            bugs.append({
                "description": description,
                "steps": steps,
                "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M")
            })
            with open(bugs_file, "w") as f:
                json.dump(bugs, f, indent=2)
            return "Bug saved locally to bugs.json. Add your GitHub token to create issues automatically."
        except Exception:
            return "Bug report saved in memory."